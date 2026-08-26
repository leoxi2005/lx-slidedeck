#!/usr/bin/env python3
"""Builds tests/fixtures/anim.pptx — a deck with real PowerPoint animation markup, used to
check the converter end to end.

python-pptx has no animation API, so the <p:timing> tree is written in by hand. It follows
the shape PowerPoint itself writes: mainSeq -> one par per click -> the effect's cTn
carrying presetClass and nodeType, with <p:spTgt spid=..><p:txEl><p:pRg st=.. end=../>.

Dev-time tool only: needs python-pptx, which the plugin itself never does.

  python3 tools/make-fixture-deck.py tests/fixtures/anim.pptx
"""
import copy, sys, pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from lxml import etree

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}


def effect_par(cid, spid, preset_class, node_type, paragraph=None, delay=0):
    """One animation effect, wrapped in the par/cTn nesting PowerPoint uses."""
    target = f'<p:spTgt spid="{spid}">'
    if paragraph is not None:
        target += f'<p:txEl><p:pRg st="{paragraph}" end="{paragraph}"/></p:txEl>'
    target += '</p:spTgt>'
    to_value = '<p:to><p:strVal val="visible"/></p:to>' if preset_class == 'entr' else \
               '<p:to><p:strVal val="hidden"/></p:to>'
    return f'''
      <p:par>
        <p:cTn id="{cid}" presetID="1" presetClass="{preset_class}" presetSubtype="0"
               fill="hold" grpId="0" nodeType="{node_type}">
          <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
          <p:childTnLst>
            <p:set>
              <p:cBhvr>
                <p:cTn id="{cid + 1}" dur="1" fill="hold">
                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                </p:cTn>
                <p:tgtEl>{target}</p:tgtEl>
                <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
              </p:cBhvr>
              {to_value}
            </p:set>
          </p:childTnLst>
        </p:cTn>
      </p:par>'''


def timing_xml(effects):
    """effects: list of (spid, presetClass, nodeType, paragraph)."""
    cid = 10
    click_groups = []
    for spid, preset, node_type, paragraph in effects:
        click_groups.append(effect_par(cid, spid, preset, node_type, paragraph))
        cid += 10
    body = "".join(click_groups)
    return f'''<p:timing xmlns:p="{NS['p']}" xmlns:a="{NS['a']}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{body}</p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''


def attach_timing(slide, effects):
    element = etree.fromstring(timing_xml(effects))
    slide.element.append(element)


def main():
    out = pathlib.Path(sys.argv[1] if len(sys.argv) > 1 else "tests/fixtures/anim.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- slide 1: a title plus four bullets that come in one click at a time -----------
    s1 = prs.slides.add_slide(blank)
    title = s1.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11), Inches(1.2))
    title.text_frame.text = "FOUR BULLETS"
    title.text_frame.paragraphs[0].font.size = Pt(44)

    body = s1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(4))
    tf = body.text_frame
    for i, line in enumerate(["Line one", "Line two", "Line three", "Line four"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(32)
    attach_timing(s1, [(body.shape_id, 'entr', 'clickEffect', n) for n in range(4)])

    # --- slide 2: one shape enters, another leaves ------------------------------------
    s2 = prs.slides.add_slide(blank)
    keep = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    keep.text_frame.text = "STAYS"
    keep.text_frame.paragraphs[0].font.size = Pt(40)

    goes = s2.shapes.add_textbox(Inches(7), Inches(1), Inches(5), Inches(2))
    goes.text_frame.text = "LEAVES"
    goes.text_frame.paragraphs[0].font.size = Pt(40)

    arrives = s2.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(2))
    arrives.text_frame.text = "ARRIVES"
    arrives.text_frame.paragraphs[0].font.size = Pt(40)

    attach_timing(s2, [
        (arrives.shape_id, 'entr', 'clickEffect', None),
        (goes.shape_id, 'exit', 'clickEffect', None),
    ])

    # --- slide 3: no animation at all, the common case --------------------------------
    s3 = prs.slides.add_slide(blank)
    plain = s3.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
    plain.text_frame.text = "NO ANIMATION"
    plain.text_frame.paragraphs[0].font.size = Pt(40)

    prs.save(str(out))
    print(f"wrote {out}")
    print("expected steps: slide1 = 5, slide2 = 3, slide3 = 1  ->  9 total")


if __name__ == "__main__":
    main()
